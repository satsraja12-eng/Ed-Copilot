"""Generate Ed-Copilot Architecture deck as a .pptx file.

Run: python scripts/generate_deck.py
Output: exports/EdCopilot_Architecture.pptx

Slides
------
 1. Title
 2. Architecture block diagram  ← full visual
 3. Plugin hook (2-file contract)
 4. DistrictRegistry startup
 5. Agent isolation
 6. YAML manifest
 7. LangGraph routing
 8. ChromaDB collections
 9. Groundedness guardrails
10. 4-step district onboarding
11. Phase roadmap
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from lxml import etree
from pptx.oxml.ns import qn

# ── Palette ──────────────────────────────────────────────────────────────────
BG_LIGHT    = RGBColor(0xF8, 0xF7, 0xF4)
BG_DARK     = RGBColor(0x1B, 0x43, 0x32)
ACCENT      = RGBColor(0x1B, 0x43, 0x32)
ACCENT_MID  = RGBColor(0x05, 0x5F, 0x3E)
ACCENT_LT   = RGBColor(0xD1, 0xFA, 0xE5)
TEXT_DARK   = RGBColor(0x0F, 0x1A, 0x13)
TEXT_LIGHT  = RGBColor(0xF8, 0xF7, 0xF4)
TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)
CODE_BG     = RGBColor(0xEC, 0xFD, 0xF5)
BORDER      = RGBColor(0xBB, 0xF7, 0xD0)
WARN_COLOR  = RGBColor(0xB4, 0x5A, 0x09)   # amber — guardrails
DB_COLOR    = RGBColor(0x1E, 0x40, 0xAF)   # blue — ChromaDB
LLM_COLOR   = RGBColor(0x4C, 0x1D, 0x95)   # violet — LLM

FONT_BODY = "Calibri"
FONT_CODE = "Courier New"

W = Inches(13.33)
H = Inches(7.5)

MARGIN_L  = Inches(0.7)
MARGIN_T  = Inches(0.6)
CONTENT_W = W - MARGIN_L * 2

RULE_Y  = Inches(1.55)
BODY_Y  = Inches(1.72)
BODY_H  = Inches(5.05)
FOOTER_Y = Inches(6.95)
FOOTER_H = Inches(0.35)
TITLE_H  = Inches(1.05)

TOTAL_SLIDES = 11


# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, fill=None, line=None, line_w_pt=0.75, radius=0):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w_pt)
    else:
        shape.line.fill.background()
    return shape


def _add_box(slide, x, y, w, h, label, sub="", fill=ACCENT,
             font_size=14, sub_size=11, bold=True):
    """Labelled block diagram box with optional sub-label."""
    shape = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = _darken(fill, 0.6)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = TEXT_LIGHT

    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = FONT_BODY
        r2.font.size = Pt(sub_size)
        r2.font.bold = False
        r2.font.color.rgb = RGBColor(0xA7, 0xF3, 0xD0) if _is_dark(fill) else TEXT_DARK

    return shape


def _darken(color: RGBColor, factor=0.7) -> RGBColor:
    return RGBColor(
        int(color[0] * factor),
        int(color[1] * factor),
        int(color[2] * factor),
    )


def _is_dark(color: RGBColor) -> bool:
    lum = 0.299 * color[0] + 0.587 * color[1] + 0.114 * color[2]
    return lum < 128


def _arrow(slide, x1, y1, x2, y2, color=TEXT_MUTED, width_pt=1.25):
    """Straight connector with arrowhead at (x2, y2)."""
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        int(x1), int(y1), int(x2), int(y2),
    )
    c.line.color.rgb = color
    c.line.width = Pt(width_pt)
    # Arrowhead at tail end (destination) via XML
    cxnSp = c._element
    spPr  = cxnSp.find(qn('p:spPr'))
    ln_el = spPr.find(qn('a:ln'))
    if ln_el is None:
        ln_el = etree.SubElement(spPr, qn('a:ln'))
    tailEnd = ln_el.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln_el, qn('a:tailEnd'))
    tailEnd.set('type', 'arrow')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    return c


def _label(slide, x, y, w, h, text, size_pt=12, bold=False,
           color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False,
           font_name=FONT_BODY):
    txb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _add_title(slide, text, dark_bg=False):
    fg = TEXT_LIGHT if dark_bg else TEXT_DARK
    txb = slide.shapes.add_textbox(MARGIN_L, MARGIN_T, CONTENT_W, TITLE_H)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = fg
    rc = TEXT_LIGHT if dark_bg else ACCENT
    _add_rect(slide, MARGIN_L, RULE_Y, CONTENT_W, Pt(2), fill=rc)


def _add_footer(slide, slide_num: int, dark_bg=False):
    fg = RGBColor(0xA0, 0xB0, 0xA8) if dark_bg else TEXT_MUTED
    _label(slide, MARGIN_L, FOOTER_Y, Inches(6), FOOTER_H,
           "Ed-Copilot Architecture", size_pt=11, color=fg)
    _label(slide, W - Inches(1.5), FOOTER_Y, Inches(1.1), FOOTER_H,
           f"{slide_num} / {TOTAL_SLIDES}", size_pt=11, color=fg,
           align=PP_ALIGN.RIGHT)


def _bullet_frame(slide, x, y, w, h):
    txb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    return tf


def _add_para(tf, text, indent=0, size_pt=20, bold=False,
              color=TEXT_DARK, font_name=FONT_BODY,
              space_before_pt=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(space_before_pt)
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_code_block(slide, x, y, w, text_lines, size_pt=14):
    line_h = Pt(size_pt) * 1.55
    block_h = int(line_h * (len(text_lines) + 1.2))
    _add_rect(slide, x, y, w, block_h, fill=CODE_BG, line=BORDER, line_w_pt=0.75)
    txb = slide.shapes.add_textbox(
        int(x + Inches(0.18)), int(y + Inches(0.12)),
        int(w - Inches(0.36)), block_h
    )
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = FONT_CODE
        run.font.size = Pt(size_pt)
        run.font.color.rgb = TEXT_DARK
    return block_h


def _add_table(slide, x, y, w, rows, col_widths, header_row=True):
    n_rows = len(rows)
    row_h  = Inches(0.42)
    tbl = slide.shapes.add_table(n_rows, len(rows[0]), int(x), int(y),
                                  int(w), row_h * n_rows).table
    total = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(w * cw / total)
    for ri in range(n_rows):
        tbl.rows[ri].height = row_h

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.text = cell_text
            run.font.name = FONT_BODY
            run.font.size = Pt(15)
            run.font.bold = (ri == 0 and header_row)
            if ri == 0 and header_row:
                run.font.color.rgb = TEXT_LIGHT
                _set_cell_bg(cell, ACCENT)
            elif "Complete" in cell_text:
                run.font.color.rgb = ACCENT
                run.font.bold = True
            elif "Pending" in cell_text:
                run.font.color.rgb = TEXT_MUTED
            else:
                run.font.color.rgb = TEXT_DARK
                if ri % 2 == 0:
                    _set_cell_bg(cell, RGBColor(0xF3, 0xF4, 0xF6))


def _set_cell_bg(cell, color: RGBColor):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing fill
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)

    _label(slide, MARGIN_L, Inches(1.8), CONTENT_W, Inches(0.4),
           "ARCHITECTURE OVERVIEW", size_pt=13, bold=True,
           color=RGBColor(0x6E, 0xE7, 0xB7))

    txb = slide.shapes.add_textbox(MARGIN_L, Inches(2.2), CONTENT_W, Inches(1.8))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Ed-Copilot"
    run.font.name = FONT_BODY
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = TEXT_LIGHT
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Multi-District K-12 Family Assistant"
    r2.font.name = FONT_BODY
    r2.font.size = Pt(28)
    r2.font.color.rgb = RGBColor(0xA7, 0xF3, 0xD0)

    _add_rect(slide, MARGIN_L, Inches(4.25), Inches(4.5), Pt(2),
              fill=RGBColor(0x6E, 0xE7, 0xB7))
    _label(slide, MARGIN_L, Inches(4.5), CONTENT_W, Inches(0.4),
           "Phase 2  ·  June 2026  ·  Engineering & Product",
           size_pt=14, color=RGBColor(0xA7, 0xF3, 0xD0))


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — Architecture Block Diagram
# ─────────────────────────────────────────────────────────────────────────────

def slide_02_diagram(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "The system routes family questions through a 5-layer plugin architecture")

    # ── Layer geometry ─────────────────────────────────────────────────────
    DY  = BODY_Y           # top of diagram area
    LH  = Inches(0.62)     # layer height
    GAP = Inches(0.32)     # vertical gap between layers
    CW  = CONTENT_W        # total content width
    LX  = MARGIN_L

    # Column layout for 3 agents / 3 DB boxes
    COL_W   = (CW - Inches(0.28)) / 3
    COL_GAP = Inches(0.14)
    col_x = [LX + i * (COL_W + COL_GAP) for i in range(3)]

    # Y positions
    y_ui      = DY
    y_orch    = y_ui   + LH + GAP
    y_agents  = y_orch + LH + GAP
    y_guard   = y_agents + LH + GAP * 0.7
    y_db      = y_guard  + LH * 0.9 + GAP * 0.5
    y_llm     = y_db     + LH + GAP * 0.9

    cx = LX + CW / 2  # horizontal center

    # ── Layer 1: User + Frontend ──────────────────────────────────────────
    # User bubble (left side)
    user_w = Inches(2.6)
    _add_box(slide, LX, y_ui, user_w, LH,
             "Family / User",
             sub="Parent · Student · Teacher",
             fill=RGBColor(0x06, 0x4E, 0x3B),
             font_size=13, sub_size=10)

    # Streamlit frontend (center+right)
    fe_x = LX + user_w + Inches(0.2)
    fe_w = CW - user_w - Inches(0.2)
    _add_box(slide, fe_x, y_ui, fe_w, LH,
             "Streamlit Frontend  (app.py)",
             sub="Persona selector  ·  District dropdown  ·  Chat interface  ·  LangSmith trace link",
             fill=ACCENT_MID, font_size=13, sub_size=10)

    # Arrow: user → frontend
    _arrow(slide, LX + user_w, y_ui + LH / 2,
           fe_x, y_ui + LH / 2, color=ACCENT)

    # ── Layer 2: Orchestrator ─────────────────────────────────────────────
    _add_box(slide, LX, y_orch, CW, LH,
             "LangGraph Orchestrator  (src/orchestrator.py)",
             sub="classify_intent node  →  routes to district node via conditional edge  |  out-of-scope short-circuits",
             fill=ACCENT, font_size=13, sub_size=10)
    # Arrow: frontend → orchestrator
    _arrow(slide, cx, y_ui + LH, cx, y_orch, color=ACCENT)

    # ── Layer 3: District Agents ──────────────────────────────────────────
    agents = [
        ("Wake County NC\nAgent", "NC Math 1/2/3\nhybrid retriever", RGBColor(0x15, 0x55, 0x3A)),
        ("Frisco ISD TX\nAgent",  "retrieve() · synthesize()\nhandle()", RGBColor(0x15, 0x55, 0x3A)),
        ("Plano ISD TX\nAgent",   "retrieve() · synthesize()\nhandle()", RGBColor(0x15, 0x55, 0x3A)),
    ]
    for i, (name, sub, col) in enumerate(agents):
        _add_box(slide, col_x[i], y_agents, COL_W, LH,
                 name, sub=sub, fill=col, font_size=12, sub_size=9)
        # Arrow: orchestrator → each agent
        ax = col_x[i] + COL_W / 2
        _arrow(slide, ax, y_orch + LH, ax, y_agents, color=ACCENT)

    # ── Layer 4a: Guardrails (TX agents only) ────────────────────────────
    guard_h = LH * 0.85
    for i in range(1, 3):
        _add_box(slide, col_x[i], y_guard, COL_W, guard_h,
                 "Guardrails",
                 sub="Safety pre-check  ·  Groundedness score",
                 fill=WARN_COLOR, font_size=11, sub_size=9)
        ax = col_x[i] + COL_W / 2
        _arrow(slide, ax, y_agents + LH, ax, y_guard, color=WARN_COLOR)

    # Wake County — label showing it skips guardrails layer
    _label(slide, col_x[0] + Inches(0.1), y_guard + Inches(0.12),
           COL_W - Inches(0.2), guard_h,
           "LangSmith eval\n(Faithfulness + Relevance)",
           size_pt=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)

    # ── Layer 4b: ChromaDB ────────────────────────────────────────────────
    dbs = [
        ("ChromaDB\nlangchain",                    "1,158 chunks\nNC Math 1/2/3"),
        ("ChromaDB\nfrisco_isd_tx__course_catalog", "7 chunks\nHS Math catalog"),
        ("ChromaDB\nplano_isd_tx__course_catalog",  "5 chunks\nHS Math catalog"),
    ]
    for i, (name, sub) in enumerate(dbs):
        _add_box(slide, col_x[i], y_db, COL_W, LH,
                 name, sub=sub, fill=DB_COLOR, font_size=11, sub_size=9)
        ax = col_x[i] + COL_W / 2
        from_y = y_guard + guard_h if i > 0 else y_agents + LH
        _arrow(slide, ax, from_y, ax, y_db, color=DB_COLOR)

    # ── Layer 5: LLM ─────────────────────────────────────────────────────
    _add_box(slide, LX, y_llm, CW, LH,
             "LLM  (Nebius · meta-llama/Llama-3.3-70B-Instruct)",
             sub="Shared across all agents  ·  Prompt includes district persona + retrieved context",
             fill=LLM_COLOR, font_size=13, sub_size=10)
    # Arrow from each DB to LLM
    for i in range(3):
        ax = col_x[i] + COL_W / 2
        _arrow(slide, ax, y_db + LH, ax, y_llm, color=LLM_COLOR)

    # ── Legend (bottom-right) ─────────────────────────────────────────────
    leg_x = W - Inches(2.8)
    leg_y = y_llm + LH + Inches(0.1)
    legend_items = [
        (ACCENT,      "Orchestration layer"),
        (DB_COLOR,    "Vector store (ChromaDB)"),
        (WARN_COLOR,  "Guardrails"),
        (LLM_COLOR,   "Language model"),
    ]
    for j, (col, lbl) in enumerate(legend_items):
        ly = leg_y + j * Inches(0.26)
        _add_rect(slide, leg_x, ly, Inches(0.22), Inches(0.2), fill=col)
        _label(slide, leg_x + Inches(0.28), ly - Inches(0.02),
               Inches(2.4), Inches(0.25), lbl, size_pt=10, color=TEXT_MUTED)

    _add_footer(slide, 2)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — Plugin hook (2-file contract)
# ─────────────────────────────────────────────────────────────────────────────

def slide_03_plugin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "Any school district can plug into Ed-Copilot by adding exactly 2 files")

    tf = _bullet_frame(slide, MARGIN_L, BODY_Y, CONTENT_W, BODY_H)
    _add_para(tf, "The 2-file contract:", size_pt=18, bold=True, color=ACCENT)
    _add_para(tf, "  config/tenants/<district-id>.yaml  —  declare the district",
              size_pt=18, font_name=FONT_CODE, space_before_pt=6)
    _add_para(tf, "  src/agents/<district_id>_agent.py  —  implement the agent",
              size_pt=18, font_name=FONT_CODE, space_before_pt=4)

    _add_para(tf, "What the orchestrator and app do automatically:", size_pt=18,
              bold=True, color=ACCENT, space_before_pt=18)
    for item in [
        "DistrictRegistry discovers the new YAML on next startup — no manual registration",
        "The sidebar district dropdown updates from the registry — no hardcoded list",
        "LangGraph adds a node for the new agent — conditional edge wired automatically",
        "ChromaDB collection names derived from district_id — no config needed",
    ]:
        _add_para(tf, f"  \u2022  {item}", size_pt=18, space_before_pt=6)

    _add_para(tf, "Current deployment: Wake County NC  ·  Frisco ISD TX  ·  Plano ISD TX",
              size_pt=17, bold=True, color=ACCENT, space_before_pt=18)
    _add_footer(slide, 3)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — DistrictRegistry startup
# ─────────────────────────────────────────────────────────────────────────────

def slide_04_registry(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "DistrictRegistry scans YAML configs at startup and wires each agent automatically")

    tf = _bullet_frame(slide, MARGIN_L, BODY_Y, CONTENT_W, BODY_H)
    _add_para(tf, "Startup sequence (src/district_registry.py):", size_pt=18,
              bold=True, color=ACCENT)
    for num, step, detail in [
        ("1.", "Scans config/tenants/*.yaml",
         "Reads all YAML files in the tenants directory on first import"),
        ("2.", "Reads agent_module from each YAML",
         "Python dotted path: src.agents.frisco_isd_tx_agent"),
        ("3.", "Imports the module and reads the  agent  variable",
         "Each agent file exports a module-level  agent = FriscoIsdAgent()  instance"),
        ("4.", "Registers  {district_id: DistrictAgent}  in an internal dict",
         "Keyed by district_id from the YAML — used for routing and display"),
        ("5.", "Orchestrator calls  registry.all_district_ids()  to build LangGraph",
         "One conditional node per district — no hardcoded district names anywhere"),
    ]:
        _add_para(tf, f"  {num}  {step}", size_pt=18, bold=True, space_before_pt=10, color=ACCENT)
        _add_para(tf, f"       {detail}", size_pt=17, space_before_pt=2)

    _add_footer(slide, 4)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — Agent isolation
# ─────────────────────────────────────────────────────────────────────────────

def slide_05_isolation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "Each district agent is isolated: shared interface, zero shared state")

    col_w = (CONTENT_W - Inches(0.3)) / 2
    lx, rx = MARGIN_L, MARGIN_L + col_w + Inches(0.3)

    for bx, label in [(lx, "DistrictAgent hooks (base class)"), (rx, "Isolation guarantees")]:
        _add_rect(slide, bx, BODY_Y, col_w, Inches(0.38), fill=ACCENT)
        _label(slide, bx + Inches(0.15), BODY_Y + Inches(0.06),
               col_w - Inches(0.3), Inches(0.3),
               label, size_pt=15, bold=True, color=TEXT_LIGHT)

    tf_l = _bullet_frame(slide, lx, BODY_Y + Inches(0.44), col_w, Inches(4.5))
    for line in [
        "retrieve(query, intent, persona)",
        "  \u2192 List[Document]",
        "",
        "synthesize(query, docs, intent, persona)",
        "  \u2192 str",
        "",
        "handle(state) \u2192 state",
        "  Full pipeline override",
        "",
        "district_id  (property)",
        "  Unique key — used for routing + DB naming",
        "",
        "supported_intents  (property)",
        "  e.g. [course_catalog, admin_policy]",
    ]:
        _add_para(tf_l, line, size_pt=15 if line.strip() else 6,
                  font_name=FONT_CODE if line.strip() and ":" not in line else FONT_BODY,
                  color=TEXT_DARK, space_before_pt=1)

    tf_r = _bullet_frame(slide, rx, BODY_Y + Inches(0.44), col_w, Inches(4.5))
    for g in [
        "Frisco ISD data never appears in a Plano ISD response",
        "Agents use isolated ChromaDB client instances — no cross-district reads",
        "Embeddings model loaded lazily per agent — independent lifecycle",
        "A crash or timeout in one agent does not affect others",
        "Each agent can use a different LLM, embedding model, or retrieval strategy",
    ]:
        _add_para(tf_r, f"  \u2022  {g}", size_pt=17, space_before_pt=12)

    _add_footer(slide, 5)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — YAML manifest
# ─────────────────────────────────────────────────────────────────────────────

def slide_06_yaml(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "The YAML manifest is the only contract between a district and the platform")

    _label(slide, MARGIN_L, BODY_Y, CONTENT_W, Inches(0.38),
           "config/tenants/frisco_isd_tx.yaml  —  complete example:",
           size_pt=17, bold=True, color=ACCENT)

    yaml_lines = [
        "district_id:   frisco_isd_tx",
        "display_name:  Frisco ISD",
        "state:         TX",
        "grade_levels:  K-12",
        "doc_types:",
        "  - course_catalog",
        "  - admin_policy",
        "agent_module:  src.agents.frisco_isd_tx_agent",
    ]
    block_h = _add_code_block(slide, MARGIN_L, BODY_Y + Inches(0.46),
                               Inches(7.2), yaml_lines, size_pt=16)

    note_y = BODY_Y + Inches(0.46) + block_h + Inches(0.22)
    tf = _bullet_frame(slide, MARGIN_L, note_y, CONTENT_W, Inches(2.2))
    _add_para(tf, "What each field drives:", size_pt=18, bold=True, color=ACCENT)
    for line in [
        "district_id     \u2192   ChromaDB collection names  (frisco_isd_tx__course_catalog)",
        "display_name    \u2192   Sidebar dropdown label in app.py",
        "agent_module    \u2192   Python import path — DistrictRegistry loads the agent class",
        "doc_types       \u2192   Intent classifier scope for this district only",
        "state           \u2192   Geographic guard — NC-specific content blocked for TX districts",
    ]:
        _add_para(tf, f"  \u2022  {line}", size_pt=16, font_name=FONT_CODE, space_before_pt=5)

    _add_footer(slide, 6)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — LangGraph routing
# ─────────────────────────────────────────────────────────────────────────────

def slide_07_langgraph(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "LangGraph routes each query to exactly one district node via intent classification")

    tf = _bullet_frame(slide, MARGIN_L, BODY_Y, CONTENT_W, BODY_H)
    _add_para(tf, "Graph execution path:", size_pt=18, bold=True, color=ACCENT)

    for num, label, desc in [
        ("1.", "classify_intent node",
         "LLM reads the user message and district context to produce: intent (course_catalog | admin_policy | out_of_scope) — no district-specific logic here"),
        ("2.", "Conditional edge",
         "Routes to agent_frisco_isd_tx, agent_plano_isd_tx, or agent_wake_county_nc based on state.district from the sidebar"),
        ("3.", "District node — handle()",
         "Calls input_is_safe() pre-check, then retrieve() to query ChromaDB, then synthesize() with the LLM and retrieved context"),
        ("4.", "Response returned",
         "State carries: response (str), context_docs (List[Document] for citation), intent_badge (displayed in UI)"),
    ]:
        _add_para(tf, f"  {num}  {label}", size_pt=19, bold=True,
                  space_before_pt=10, color=ACCENT)
        _add_para(tf, f"       {desc}", size_pt=17, space_before_pt=2)

    _add_para(tf, "Out-of-scope queries short-circuit at the conditional edge — no agent called",
              size_pt=17, italic=True, color=TEXT_MUTED, space_before_pt=14)
    _add_footer(slide, 7)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — ChromaDB collections
# ─────────────────────────────────────────────────────────────────────────────

def slide_08_chromadb(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "Frisco and Plano agents retrieve from district-isolated ChromaDB collections")

    _label(slide, MARGIN_L, BODY_Y, CONTENT_W, Inches(0.38),
           "Collection naming:   {district_id}__{doc_type}   —   enforces zero cross-district data leakage",
           size_pt=17, bold=True, color=ACCENT)

    rows = [
        ["Collection", "Chunks", "District", "Content"],
        ["frisco_isd_tx__course_catalog",   "7",     "Frisco ISD TX",   "AP Calc AB/BC, Geometry, Algebra I/II, Precalculus, AP Stats"],
        ["plano_isd_tx__course_catalog",    "5",     "Plano ISD TX",    "Algebra I through AP Calculus AB"],
        ["langchain",                        "1,158", "Wake County NC",  "NC Math 1, 2, 3 curriculum + standards"],
        ["frisco_isd_tx__admin_policy",     "0",     "Frisco ISD TX",   "Admin PDFs — ingestion pending (Phase 3)"],
        ["plano_isd_tx__admin_policy",      "0",     "Plano ISD TX",    "Admin PDFs — ingestion pending (Phase 3)"],
    ]
    _add_table(slide, MARGIN_L, BODY_Y + Inches(0.46),
               CONTENT_W, rows, col_widths=[4, 1, 2.2, 5.5])

    footer_y = BODY_Y + Inches(0.46) + Inches(0.42 * len(rows)) + Inches(0.18)
    _label(slide, MARGIN_L, footer_y, CONTENT_W, Inches(0.55),
           "Embeddings: BAAI/bge-small-en-v1.5 (HuggingFace) — consistent across all districts.\n"
           "Ingestion: seed JSON (offline) or --crawl flag to fetch from live district websites.",
           size_pt=13, color=TEXT_MUTED)

    _add_footer(slide, 8)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — Groundedness guardrails
# ─────────────────────────────────────────────────────────────────────────────

def slide_09_guardrails(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "Groundedness guardrails prevent hallucinated and unsafe responses before delivery")

    tf = _bullet_frame(slide, MARGIN_L, BODY_Y, CONTENT_W, BODY_H)
    _add_para(tf, "Two-layer guard — src/guardrails/groundedness.py:", size_pt=18,
              bold=True, color=ACCENT)

    _add_para(tf, "  Layer 1  —  Safety pre-check  (input_is_safe)", size_pt=19,
              bold=True, space_before_pt=14)
    for line in [
        "Fires synchronously before any retrieval — zero latency cost",
        "Blocks PII-elicitation, student-privacy probes, and data harvesting prompts",
        "Returns a refusal string; handle() exits without calling retrieve() or synthesize()",
    ]:
        _add_para(tf, f"       \u2022  {line}", size_pt=17, space_before_pt=4)

    _add_para(tf, "  Layer 2  —  Groundedness score  (post-synthesis)", size_pt=19,
              bold=True, space_before_pt=14)
    for line in [
        "Lexical overlap: fraction of non-stop answer tokens that appear in retrieved context",
        "Score < 0.25: warning logged server-side; response still delivered (non-blocking at POC stage)",
        "Upgrade path: NLI cross-encoder entailment or LLM-as-judge for production blocking",
    ]:
        _add_para(tf, f"       \u2022  {line}", size_pt=17, space_before_pt=4)

    _add_para(tf, "  Wake County NC: LangSmith Faithfulness + Relevance eval  (LLM-as-judge, 15 Q&A pairs)",
              size_pt=17, bold=True, color=ACCENT, space_before_pt=14)

    _add_footer(slide, 9)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — 4-step onboarding
# ─────────────────────────────────────────────────────────────────────────────

def slide_10_onboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "Adding a 4th district requires one YAML file, one Python class, and an ingestion run")

    _label(slide, MARGIN_L, BODY_Y, CONTENT_W, Inches(0.38),
           "Example: onboarding Richardson ISD TX", size_pt=18, bold=True, color=ACCENT)

    steps = [
        ("Step 1", "config/tenants/richardson_isd_tx.yaml",
         "8 fields — district_id, display_name, state, grade_levels, doc_types, agent_module"),
        ("Step 2", "src/agents/richardson_isd_tx_agent.py",
         "Subclass DistrictAgent, implement retrieve() + synthesize() — copy Frisco agent as template"),
        ("Step 3", "python src/ingestion/richardson_ingestion.py",
         "Seed JSON offline; use --crawl flag to fetch live from district website. Populates ChromaDB."),
        ("Step 4", "Restart app",
         "Richardson ISD appears in the district dropdown automatically — zero other changes"),
    ]

    step_y = BODY_Y + Inches(0.46)
    step_h = Inches(1.06)
    for i, (label, code, desc) in enumerate(steps):
        by = step_y + i * (step_h + Inches(0.08))
        _add_rect(slide, MARGIN_L, by, Inches(0.45), step_h, fill=ACCENT)
        _label(slide, MARGIN_L + Inches(0.05), by + Inches(0.26),
               Inches(0.35), Inches(0.5),
               str(i + 1), size_pt=22, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        _label(slide, MARGIN_L + Inches(0.6), by + Inches(0.08),
               CONTENT_W - Inches(0.65), Inches(0.42),
               code, size_pt=16, bold=True, color=TEXT_DARK, font_name=FONT_CODE)
        _label(slide, MARGIN_L + Inches(0.6), by + Inches(0.52),
               CONTENT_W - Inches(0.65), Inches(0.46),
               desc, size_pt=16, color=TEXT_DARK)

    note_y = step_y + 4 * (step_h + Inches(0.08)) + Inches(0.05)
    _label(slide, MARGIN_L, note_y, CONTENT_W, Inches(0.38),
           "No changes to: orchestrator.py  ·  app.py  ·  district_registry.py  ·  any other district's code",
           size_pt=14, color=TEXT_MUTED, italic=True)

    _add_footer(slide, 10)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 11 — Phase roadmap
# ─────────────────────────────────────────────────────────────────────────────

def slide_11_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_LIGHT)
    _add_title(slide, "Phases 1 and 2 are complete; Phases 3-5 deliver user profiles, evals, and citations")

    rows = [
        ["Phase", "Deliverable", "Status"],
        ["Phase 1", "Plugin hook system — DistrictAgent base class, DistrictRegistry, YAML auto-discovery", "Complete"],
        ["Phase 2", "TX ingestion & live retrieval — crawlers, pipeline, groundedness, Frisco + Plano agents live", "Complete"],
        ["Phase 3", "User profiles — SQLite session registry, auto-routing from saved district + role", "Pending"],
        ["Phase 4", "Evaluation — Frisco/Plano LangSmith gold Q&A pairs, Faithfulness + Relevance scores", "Pending"],
        ["Phase 5", "Citations — source URLs surfaced inline in every response", "Pending"],
    ]
    _add_table(slide, MARGIN_L, BODY_Y, CONTENT_W, rows,
               col_widths=[1.5, 8.5, 2])
    _add_footer(slide, 11)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_diagram(prs)
    slide_03_plugin(prs)
    slide_04_registry(prs)
    slide_05_isolation(prs)
    slide_06_yaml(prs)
    slide_07_langgraph(prs)
    slide_08_chromadb(prs)
    slide_09_guardrails(prs)
    slide_10_onboard(prs)
    slide_11_roadmap(prs)

    os.makedirs("exports", exist_ok=True)
    out = "exports/EdCopilot_Architecture.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
    return out


if __name__ == "__main__":
    build()
